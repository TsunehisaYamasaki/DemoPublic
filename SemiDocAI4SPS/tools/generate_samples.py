#!/usr/bin/env python3
"""
半導体チーム向けサンプル入力ファイル生成スクリプト (SemiDocAI4SPS 版)
================================================================
DI4Storage の generate_samples.py を参考に、SharePoint Online へアップロードする
サンプルドキュメントを自動生成する。

生成されるファイル:
  1. PNG  - WaveDrom タイミングダイアグラム画像
  2. DOCX - SoC 設計仕様書
  3. XLSX - ウェーハテスト結果 (歩留まりデータ)
  4. PPTX - 設計レビュー資料
  5. PDF  - 品質検査報告書

Usage:
    pip install python-docx openpyxl python-pptx Pillow reportlab
    python tools/generate_samples.py
"""
from pathlib import Path
import json, math, random, struct, zlib

# 出力先ディレクトリ
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "sample_inputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ═══════════════════════════════════════════════════════════
# 1. PNG - WaveDrom タイミングダイアグラム画像 (Pure Python)
# ═══════════════════════════════════════════════════════════
def generate_wavedrom_image():
    """SoC バスタイミングのダイアグラムを Pure Python PNG として生成"""

    wavedrom_json = {
        "signal": [
            {"name": "CLK",        "wave": "P.........",  "period": 1},
            {"name": "HADDR[31:0]","wave": "x.3.4.5.x.", "data": ["A0","A1","A2"]},
            {"name": "HWDATA",     "wave": "x...3.4.x.", "data": ["D0","D1"]},
            {"name": "HWRITE",     "wave": "0..1....0."},
            {"name": "HREADY",     "wave": "1..0.1..1."},
            {"name": "HRESP",      "wave": "0........0"},
            {"name": "HRDATA",     "wave": "x.......3.", "data": ["R0"]},
        ],
        "head": {"text": "SoC AHB-Lite Bus Timing (7nm FinFET)", "tick": 0},
        "foot": {"text": "Document: SSS-TIMING-2026-001 | Confidential"},
    }

    width, height = 900, 420
    bg = (255, 255, 255)

    pixels = []
    for _ in range(height):
        row = []
        for _ in range(width):
            row.append(bg)
        pixels.append(row)

    def draw_rect(px, x0, y0, w, h, color):
        for yy in range(max(0, y0), min(height, y0 + h)):
            for xx in range(max(0, x0), min(width, x0 + w)):
                px[yy][xx] = color

    def draw_hline(px, x0, x1, y, color):
        if 0 <= y < height:
            for xx in range(max(0, x0), min(width, x1)):
                px[y][xx] = color

    def draw_vline(px, x, y0, y1, color):
        if 0 <= x < width:
            for yy in range(max(0, y0), min(height, y1)):
                px[yy][x] = color

    # 背景ヘッダ
    draw_rect(pixels, 0, 0, width, 40, (26, 35, 126))  # 紺色ヘッダ

    # 信号描画
    signals = wavedrom_json["signal"]
    y_start = 60
    row_h = 45
    x_label_w = 140
    step_w = 70

    colors = {
        "sig_high": (0, 128, 0),
        "sig_low": (0, 128, 0),
        "clk": (0, 0, 200),
        "data": (200, 100, 0),
        "grid": (220, 220, 220),
        "text_area": (240, 245, 255),
    }

    for idx, sig in enumerate(signals):
        y_base = y_start + idx * row_h
        wave = sig.get("wave", "")

        # ラベル領域
        draw_rect(pixels, 0, y_base, x_label_w, row_h - 5, (245, 245, 250))

        # 波形描画
        for step_i, ch in enumerate(wave):
            x0 = x_label_w + step_i * step_w
            x1 = x0 + step_w
            y_mid = y_base + row_h // 2 - 2

            # グリッド
            draw_vline(pixels, x0, y_base, y_base + row_h - 5, colors["grid"])

            if ch == 'P':  # クロック High
                draw_hline(pixels, x0, x0 + step_w // 2, y_base + 8, colors["clk"])
                draw_vline(pixels, x0 + step_w // 2, y_base + 8, y_base + row_h - 8, colors["clk"])
                draw_hline(pixels, x0 + step_w // 2, x1, y_base + row_h - 8, colors["clk"])
            elif ch == '1':
                draw_hline(pixels, x0, x1, y_base + 8, colors["sig_high"])
            elif ch == '0':
                draw_hline(pixels, x0, x1, y_base + row_h - 8, colors["sig_low"])
            elif ch in ('3', '4', '5'):  # data
                draw_rect(pixels, x0 + 2, y_base + 6, step_w - 4, row_h - 16, colors["text_area"])
                draw_hline(pixels, x0, x1, y_base + 6, colors["data"])
                draw_hline(pixels, x0, x1, y_base + row_h - 10, colors["data"])
            elif ch == 'x':  # unknown
                for xx in range(x0, min(x1, width), 6):
                    draw_hline(pixels, xx, min(xx + 3, width), y_mid, (200, 200, 200))
            elif ch == '.':
                pass  # 前の状態を継続 (簡略化)

    # フッタ
    draw_rect(pixels, 0, height - 30, width, 30, (240, 240, 240))

    # --- PNG エンコード (Pure Python) ---
    def encode_png(px, w, h):
        def chunk(ctype, data):
            c = ctype + data
            crc = zlib.crc32(c) & 0xFFFFFFFF
            return struct.pack(">I", len(data)) + c + struct.pack(">I", crc)

        raw = b""
        for row in px:
            raw += b"\x00"
            for r, g, b in row:
                raw += struct.pack("BBB", r, g, b)

        sig = b"\x89PNG\r\n\x1a\n"
        ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
        return sig + chunk(b"IHDR", ihdr) + chunk(b"IDAT", zlib.compress(raw, 9)) + chunk(b"IEND", b"")

    out_path = OUTPUT_DIR / "design_wavedrom_soc_timing.png"
    png_data = encode_png(pixels, width, height)
    out_path.write_bytes(png_data)

    # JSON も保存
    json_path = OUTPUT_DIR / "design_wavedrom_soc_timing.json"
    json_path.write_text(json.dumps(wavedrom_json, indent=2, ensure_ascii=False), encoding="utf-8")

    print(f"[OK] WaveDrom PNG: {out_path}")
    print(f"[OK] WaveDrom JSON: {json_path}")


# ═══════════════════════════════════════════════════════════
# 2. DOCX - 設計部門 SoC 仕様書
# ═══════════════════════════════════════════════════════════
def generate_word_document():
    """SoC 設計仕様書を Word ドキュメント (DOCX) として生成"""
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()

    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(10.5)

    # ─── 表紙 ───
    for _ in range(6):
        doc.add_paragraph()

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("SSS-7NM-A1 SoC Design Specification")
    run.font.size = Pt(26)
    run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    run.bold = True

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("TSMC 7nm FinFET Process")
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    for _ in range(4):
        doc.add_paragraph()

    info_items = [
        ("Document No.", "SSS-SPEC-2026-001"),
        ("Revision", "Rev 2.3"),
        ("Author", "Design Dept. / Sato"),
        ("Reviewer", "Nakashima"),
        ("Date", "2026-02-15"),
        ("Classification", "Confidential"),
    ]
    info_table = doc.add_table(rows=len(info_items), cols=2)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, (k, v) in enumerate(info_items):
        info_table.rows[i].cells[0].text = k
        info_table.rows[i].cells[1].text = v
        for cell in info_table.rows[i].cells:
            for p in cell.paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # ─── 目次 (テキストベース) ───
    doc.add_heading("Table of Contents", level=1)
    toc_items = [
        "1. Overview", "2. Architecture", "3. Block Descriptions",
        "4. Clock & Reset", "5. Memory Map", "6. Electrical Characteristics",
        "7. Package", "8. Revision History",
    ]
    for item in toc_items:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(2)

    doc.add_page_break()

    # ─── 1. 概要 ───
    doc.add_heading("1. Overview", level=1)
    doc.add_paragraph(
        "The SSS-7NM-A1 is a high-performance SoC designed for next-generation "
        "semiconductor manufacturing equipment control. It integrates a quad-core "
        "CPU cluster, GPU, DSP, and dedicated hardware accelerators on TSMC's "
        "7nm FinFET (N7) process."
    )

    spec_data = [
        ["Item",           "Specification"],
        ["Process",        "TSMC 7nm FinFET (N7)"],
        ["Die Size",       "8.5mm x 9.2mm (78.2 mm²)"],
        ["Transistor Count","~1.2 Billion"],
        ["Max Frequency",  "1.2 GHz"],
        ["Supply Voltage", "0.75V (Core) / 1.8V (IO)"],
        ["Power",          "2.8W TDP (Typical)"],
        ["Package",        "FCBGA 23x23mm, 784 balls, 0.8mm pitch"],
        ["Operating Temp", "-40°C to 105°C (Industrial)"],
    ]
    table = doc.add_table(rows=len(spec_data), cols=2)
    table.style = "Light Grid Accent 1"
    for i, row_data in enumerate(spec_data):
        for j, val in enumerate(row_data):
            table.rows[i].cells[j].text = val

    # ─── 2. アーキテクチャ ───
    doc.add_heading("2. Architecture", level=1)
    doc.add_paragraph(
        "The SoC architecture consists of the following major subsystems connected "
        "via a high-performance AHB-Lite bus fabric:"
    )

    blocks = [
        ("CPU Cluster", "Quad-core ARM Cortex-A55-based, 32KB L1I + 32KB L1D per core, 512KB shared L2"),
        ("GPU", "Mali-based 2-shader core, OpenCL 2.0 support, 25.6 GFLOPS peak"),
        ("DSP Subsystem", "Dual VLIW DSP cores, 256-bit SIMD, 19.2 GOPS peak performance"),
        ("DMA Controller", "8-channel DMA with scatter-gather, priority arbitration"),
        ("Memory Controller", "LPDDR4X-4266, 32-bit, ECC support, max 4GB"),
        ("Bus Fabric", "AHB-Lite multi-layer interconnect, 128-bit data width"),
        ("Peripheral Hub", "UART x4, SPI x3, I2C x4, GPIO x32, CAN-FD x2, Ethernet 1G"),
        ("Security Engine", "AES-256, SHA-3, RSA-4096, TRNG, Secure Boot, TrustZone"),
        ("PLL / Clock Gen", "3x PLL, frequency range 100MHz-1.5GHz, spread spectrum"),
        ("Power Management", "DVFS support, 4 power domains, retention mode"),
    ]
    for name, desc in blocks:
        doc.add_heading(name, level=3)
        doc.add_paragraph(desc)

    # ─── 3. メモリマップ ───
    doc.add_heading("5. Memory Map", level=1)

    mmap = [
        ["Address Range",           "Size",   "Block"],
        ["0x0000_0000 - 0x0007_FFFF", "512KB", "Boot ROM"],
        ["0x1000_0000 - 0x1007_FFFF", "512KB", "SRAM (TCM)"],
        ["0x2000_0000 - 0x200F_FFFF", "1MB",   "Peripheral registers"],
        ["0x4000_0000 - 0x7FFF_FFFF", "1GB",   "DRAM (LPDDR4X)"],
        ["0x8000_0000 - 0x8000_FFFF", "64KB",  "GPU registers"],
        ["0x9000_0000 - 0x9000_FFFF", "64KB",  "DSP registers"],
        ["0xA000_0000 - 0xA000_0FFF", "4KB",   "Security Engine"],
        ["0xE000_0000 - 0xE00F_FFFF", "1MB",   "Debug / CoreSight"],
    ]
    table = doc.add_table(rows=len(mmap), cols=3)
    table.style = "Light Grid Accent 1"
    for i, row_data in enumerate(mmap):
        for j, val in enumerate(row_data):
            table.rows[i].cells[j].text = val

    # ─── 6. 電気特性 ───
    doc.add_heading("6. Electrical Characteristics", level=1)

    elec = [
        ["Parameter",             "Min",   "Typ",   "Max",   "Unit"],
        ["Core Voltage (VDD)",    "0.675", "0.750", "0.825", "V"],
        ["IO Voltage (VDDIO)",    "1.62",  "1.80",  "1.98",  "V"],
        ["Dynamic Power @1.2GHz", "-",     "1.92",  "2.50",  "W"],
        ["Leakage Power @25°C",   "-",     "0.88",  "1.20",  "W"],
        ["Iddq @0.75V",           "-",     "2.1",   "3.5",   "mA"],
        ["Fmax",                  "1.10",  "1.20",  "-",     "GHz"],
        ["IO Drive Strength",     "2",     "4",     "8",     "mA"],
    ]
    table = doc.add_table(rows=len(elec), cols=5)
    table.style = "Light Grid Accent 1"
    for i, row_data in enumerate(elec):
        for j, val in enumerate(row_data):
            table.rows[i].cells[j].text = val

    # ─── 改訂履歴 ───
    doc.add_heading("8. Revision History", level=1)
    rev_data = [
        ["Rev", "Date",       "Author",     "Description"],
        ["1.0", "2025-08-01", "Sato",       "Initial release"],
        ["1.1", "2025-10-15", "Sato",       "Added DSP block specification"],
        ["2.0", "2025-12-20", "Nakashima",  "Major update: Memory map revision, Security Engine added"],
        ["2.1", "2026-01-10", "Sato",       "Updated electrical characteristics"],
        ["2.2", "2026-01-28", "Sato",       "Added CAN-FD peripheral"],
        ["2.3", "2026-02-15", "Sato",       "Clock tree & PLL specification update"],
    ]
    table = doc.add_table(rows=len(rev_data), cols=4)
    table.style = "Light Grid Accent 1"
    for i, row_data in enumerate(rev_data):
        for j, val in enumerate(row_data):
            table.rows[i].cells[j].text = val

    out_path = OUTPUT_DIR / "design_soc_specification.docx"
    doc.save(str(out_path))
    print(f"[OK] Word SoC Specification: {out_path}")


# ═══════════════════════════════════════════════════════════
# 3. XLSX - 製造部門 ウェーハテスト結果
# ═══════════════════════════════════════════════════════════
def generate_excel_document():
    """ウェーハテスト結果を Excel ファイルとして生成"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    random.seed(42)

    # ── Sheet 1: Wafer Test Summary ──
    ws1 = wb.active
    ws1.title = "Wafer Test Summary"

    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1A237E", end_color="1A237E", fill_type="solid")
    pass_fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
    fail_fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")
    border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    ws1.append(["Semiconductor Wafer Test Results - Lot: LOT-2026-0216-A"])
    ws1.merge_cells("A1:H1")
    ws1["A1"].font = Font(name="Calibri", size=14, bold=True, color="1A237E")
    ws1.append(["Process: TSMC 7nm FinFET (N7) | Product: SSS-7NM-A1 | Date: 2026-02-16"])
    ws1.merge_cells("A2:H2")
    ws1.append([])

    headers = ["Wafer ID", "Total Dies", "Good Dies", "Defective", "Yield(%)",
               "Avg Iddq(mA)", "Avg Fmax(GHz)", "Result"]
    ws1.append(headers)
    for col_i, _ in enumerate(headers, 1):
        cell = ws1.cell(row=4, column=col_i)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
        cell.border = border

    total_dies = 892
    for i in range(1, 26):
        good = random.randint(790, 870)
        bad = total_dies - good
        yld = round(good / total_dies * 100, 1)
        iddq = round(random.uniform(1.5, 3.8), 2)
        fmax = round(random.uniform(1.08, 1.30), 2)
        result = "PASS" if yld >= 90 and iddq < 3.5 and fmax >= 1.1 else "FAIL"

        row = [f"W{i:03d}", total_dies, good, bad, yld, iddq, fmax, result]
        ws1.append(row)
        row_num = i + 4
        for col_i in range(1, 9):
            cell = ws1.cell(row=row_num, column=col_i)
            cell.alignment = Alignment(horizontal="center")
            cell.border = border
        result_cell = ws1.cell(row=row_num, column=8)
        result_cell.fill = pass_fill if result == "PASS" else fail_fill

    for col_i in range(1, 9):
        ws1.column_dimensions[get_column_letter(col_i)].width = 16

    # ── Sheet 2: PCM Data ──
    ws2 = wb.create_sheet("PCM Data")

    ws2.append(["Parametric Test Results (PCM) - Lot: LOT-2026-0216-A"])
    ws2.merge_cells("A1:G1")
    ws2["A1"].font = Font(name="Calibri", size=14, bold=True, color="1A237E")
    ws2.append([])

    pcm_headers = ["Parameter", "Lower Spec", "Upper Spec", "Unit", "Measured", "Margin(%)", "Result"]
    ws2.append(pcm_headers)
    for col_i, _ in enumerate(pcm_headers, 1):
        cell = ws2.cell(row=3, column=col_i)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
        cell.border = border

    pcm_data = [
        ["Vth (NMOS)",    0.28, 0.38, "V",      0.33],
        ["Vth (PMOS)",   -0.38,-0.28, "V",     -0.34],
        ["Idsat (NMOS)",  750,  None,  "uA/um",  820],
        ["Idsat (PMOS)",  350,  None,  "uA/um",  392],
        ["Ioff (NMOS)",   None, 10,    "nA/um",  4.2],
        ["Ioff (PMOS)",   None, 8,     "nA/um",  3.1],
        ["Tox",           1.8,  2.2,   "nm",     2.01],
        ["Rsh (Poly)",    180,  220,   "ohm/sq", 198],
        ["Rsh (Metal1)",  0.04, 0.06,  "ohm/sq", 0.049],
        ["Rsh (Metal2)",  0.03, 0.05,  "ohm/sq", 0.041],
        ["Via Resistance", 0.8, 1.5,   "ohm",    1.1],
        ["Contact Resist", 20,  50,    "ohm",    32],
    ]
    for row_data in pcm_data:
        name, lo, hi, unit, measured = row_data
        if lo is not None and hi is not None:
            margin = round(min(abs(measured - lo), abs(hi - measured)) / abs(hi - lo) * 100, 1)
        elif lo is not None:
            margin = round((measured - lo) / abs(lo) * 100, 1)
        else:
            margin = round((hi - measured) / abs(hi) * 100, 1)
        result = "PASS"
        ws2.append([name, lo if lo else "-", hi if hi else "-", unit, measured, margin, result])
        rn = ws2.max_row
        for ci in range(1, 8):
            cell = ws2.cell(row=rn, column=ci)
            cell.alignment = Alignment(horizontal="center")
            cell.border = border
        ws2.cell(row=rn, column=7).fill = pass_fill

    for col_i in range(1, 8):
        ws2.column_dimensions[get_column_letter(col_i)].width = 16

    out_path = OUTPUT_DIR / "manufacturing_wafer_test_results.xlsx"
    wb.save(str(out_path))
    print(f"[OK] Excel Wafer Test: {out_path}")


# ═══════════════════════════════════════════════════════════
# 4. PPTX - 設計レビュー資料
# ═══════════════════════════════════════════════════════════
def generate_pptx_document():
    """設計レビュー資料を PowerPoint として生成"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── スライド1: 表紙 ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide1.shapes.add_textbox(Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    tf = bg.text_frame
    p = tf.paragraphs[0]
    p.text = ""

    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SSS-7NM-A1 SoC Design Review"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    p2 = tf.add_paragraph()
    p2.text = "TSMC 7nm FinFET Process | Phase 2 Review"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    info_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    for line in ["Date: 2026-02-15", "Presenter: Design Dept. Sato / Nakashima",
                 "Classification: Confidential"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── スライド2: 目次 ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title2.text_frame
    p = tf.paragraphs[0]
    p.text = "Agenda"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    agenda_items = [
        "1. Project Status Overview",
        "2. Architecture & Block Diagram",
        "3. RTL Design Progress",
        "4. Synthesis Results",
        "5. Issues & Countermeasures",
    ]
    agenda_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf = agenda_box.text_frame
    for item in agenda_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.space_after = Pt(18)

    # ── スライド3: プロジェクト状況 ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title3.text_frame
    p = tf.paragraphs[0]
    p.text = "1. Project Status Overview"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    status_items = [
        ("CPU Core RTL",      "Complete",      "100%"),
        ("GPU Block RTL",     "Complete",      "100%"),
        ("DSP Subsystem RTL", "In Progress",   "85%"),
        ("Bus Fabric RTL",    "Complete",      "100%"),
        ("DMA Controller",    "In Progress",   "90%"),
        ("Memory Controller", "Complete",      "100%"),
        ("PLL / Clock Gen",   "In Progress",   "70%"),
        ("Security Engine",   "Complete",      "100%"),
        ("Verification",      "In Progress",   "78%"),
    ]
    # テーブル
    rows, cols = len(status_items) + 1, 3
    table = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(10), Inches(5)).table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)
    for ci, h in enumerate(["Block", "Status", "Progress"]):
        cell = table.cell(0, ci)
        cell.text = h
        for pg in cell.text_frame.paragraphs:
            pg.font.size = Pt(16)
            pg.font.bold = True
    for ri, (block, status, pct) in enumerate(status_items, 1):
        table.cell(ri, 0).text = block
        table.cell(ri, 1).text = status
        table.cell(ri, 2).text = pct

    # ── スライド4: 論理合成結果 ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title4.text_frame
    p = tf.paragraphs[0]
    p.text = "4. Synthesis Results"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    synth_items = [
        "■ Target Frequency: 1.2 GHz (Clock Period: 0.833ns)",
        "■ Worst Negative Slack (WNS): +0.05ns → Constraint Met",
        "■ Total Negative Slack (TNS): 0ns",
        "■ Area: 2.85 mm² (Gate Count: ~8.5M gates)",
        "■ Dynamic Power: 1.92W @ 1.2GHz, 0.75V",
        "■ Leakage Power: 0.88W",
        "■ Critical Path: CPU Core → Bus Fabric → DMA (3-stage pipeline)",
    ]
    synth_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5))
    tf = synth_box.text_frame
    for item in synth_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    # ── スライド5: 課題と対策 ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title5.text_frame
    p = tf.paragraphs[0]
    p.text = "5. Issues & Countermeasures"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    issues = [
        ("Issue 1", "PLL/Clock Gen block RTL delay",
         "Action: Additional resources allocated. RTL completion by 2/28"),
        ("Issue 2", "DMA Controller verification coverage insufficient (75%)",
         "Action: Add random verification scenarios. Target 95%+"),
        ("Issue 3", "Bus Fabric timing margin is narrow (+0.05ns)",
         "Action: Considering additional pipeline stage"),
    ]
    issue_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5))
    tf = issue_box.text_frame
    for title_text, desc, action in issues:
        p = tf.add_paragraph()
        p.text = f"● {title_text}: {desc}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_before = Pt(16)

        p2 = tf.add_paragraph()
        p2.text = f"  → {action}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0x00, 0x80, 0x00)

    out_path = OUTPUT_DIR / "design_review_presentation.pptx"
    prs.save(str(out_path))
    print(f"[OK] PowerPoint Review: {out_path}")


# ═══════════════════════════════════════════════════════════
# 5. PDF - 品質検査報告書
# ═══════════════════════════════════════════════════════════
def generate_pdf_document():
    """品質検査報告書を PDF として生成"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import HexColor, black, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, HRFlowable,
    )
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    out_path = OUTPUT_DIR / "manufacturing_quality_report.pdf"

    doc = SimpleDocTemplate(
        str(out_path),
        pagesize=A4,
        topMargin=25 * mm,
        bottomMargin=20 * mm,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
    )

    # フォント登録
    try:
        pdfmetrics.registerFont(TTFont("MSGothic", "msgothic.ttc"))
        JP_FONT = "MSGothic"
    except Exception:
        JP_FONT = "Helvetica"

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        "TitleJP", fontName=JP_FONT, fontSize=20, leading=26,
        alignment=TA_CENTER, spaceAfter=6 * mm,
    ))
    styles.add(ParagraphStyle(
        "SubtitleJP", fontName=JP_FONT, fontSize=12, leading=16,
        alignment=TA_CENTER, textColor=HexColor("#555555"), spaceAfter=4 * mm,
    ))
    styles.add(ParagraphStyle(
        "HeadingJP", fontName=JP_FONT, fontSize=14, leading=18,
        spaceBefore=6 * mm, spaceAfter=3 * mm, textColor=HexColor("#1A237E"),
    ))
    styles.add(ParagraphStyle(
        "BodyJP", fontName=JP_FONT, fontSize=10, leading=15, spaceAfter=2 * mm,
    ))
    styles.add(ParagraphStyle(
        "SmallJP", fontName=JP_FONT, fontSize=8, leading=11,
        textColor=HexColor("#888888"),
    ))

    elements = []

    # ─── 表紙 ───
    elements.append(Spacer(1, 40 * mm))
    elements.append(Paragraph("Semiconductor Quality Inspection Report", styles["TitleJP"]))
    elements.append(Paragraph("Quality Inspection Report", styles["SubtitleJP"]))
    elements.append(Spacer(1, 10 * mm))
    elements.append(HRFlowable(width="80%", thickness=1, color=HexColor("#1A237E")))
    elements.append(Spacer(1, 10 * mm))

    cover_data = [
        ["Document No.", "SSS-QA-2026-0216-001"],
        ["Lot No.", "LOT-2026-0216-A"],
        ["Process", "TSMC 7nm FinFET (N7)"],
        ["Product", "SSS-7NM-A1 (SoC)"],
        ["Inspector", "Sato / Nakashima"],
        ["Date", "2026-02-16"],
        ["Classification", "Confidential"],
    ]
    cover_table = Table(cover_data, colWidths=[50 * mm, 80 * mm])
    cover_table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), JP_FONT),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("TEXTCOLOR", (0, 0), (0, -1), HexColor("#1A237E")),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("LINEBELOW", (0, 0), (-1, -2), 0.5, HexColor("#CCCCCC")),
        ("LINEBELOW", (0, -1), (-1, -1), 1, HexColor("#1A237E")),
    ]))
    elements.append(cover_table)

    elements.append(PageBreak())

    # ─── 1. Inspection Overview ───
    elements.append(Paragraph("1. Inspection Overview", styles["HeadingJP"]))
    elements.append(Paragraph(
        "This report documents the quality inspection results for Lot LOT-2026-0216-A "
        "of the next-generation SoC 'SSS-7NM-A1'. The inspection covers wafer-level testing, "
        "parametric measurements (PCM), appearance inspection, and reliability testing. "
        "Target yield is 92% or higher.",
        styles["BodyJP"],
    ))
    elements.append(Paragraph(
        "Wafer diameter: 300mm, dies per wafer: 892, die size: 8.5mm x 9.2mm (78.2mm2).",
        styles["BodyJP"],
    ))

    # ─── 2. Wafer Test Results ───
    elements.append(Paragraph("2. Wafer Test Results Summary", styles["HeadingJP"]))

    random.seed(42)
    wt_header = ["Wafer ID", "Good Dies", "Defective", "Yield(%)", "Iddq(mA)", "Fmax(GHz)", "Result"]
    wt_rows = [wt_header]
    pass_count = 0
    for i in range(1, 16):
        good = random.randint(790, 865)
        bad = 892 - good
        yld = round(good / 892 * 100, 1)
        iddq = round(random.uniform(1.5, 3.8), 2)
        fmax = round(random.uniform(1.08, 1.30), 2)
        result = "PASS" if yld >= 90 and iddq < 3.5 and fmax >= 1.1 else "FAIL"
        if result == "PASS":
            pass_count += 1
        wt_rows.append([f"W{i:03d}", str(good), str(bad), str(yld), str(iddq), str(fmax), result])

    wt_table = Table(wt_rows, colWidths=[20*mm, 20*mm, 20*mm, 20*mm, 20*mm, 22*mm, 16*mm])
    wt_style = [
        ("FONTNAME", (0, 0), (-1, -1), JP_FONT),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1A237E")),
        ("TEXTCOLOR", (0, 0), (-1, 0), white),
        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]
    for row_idx in range(1, len(wt_rows)):
        if wt_rows[row_idx][-1] == "PASS":
            wt_style.append(("BACKGROUND", (-1, row_idx), (-1, row_idx), HexColor("#C6EFCE")))
        else:
            wt_style.append(("BACKGROUND", (-1, row_idx), (-1, row_idx), HexColor("#FFC7CE")))
    wt_table.setStyle(TableStyle(wt_style))
    elements.append(wt_table)
    elements.append(Spacer(1, 3 * mm))
    elements.append(Paragraph(
        f"Summary: PASS {pass_count} wafers / FAIL {15 - pass_count} wafers, "
        f"Average Yield {round(sum(float(r[3]) for r in wt_rows[1:]) / 15, 1)}%",
        styles["BodyJP"],
    ))

    # ─── 3. PCM ───
    elements.append(Paragraph("3. Parametric Test Results (PCM)", styles["HeadingJP"]))
    pcm_header = ["Parameter", "Lower Limit", "Upper Limit", "Unit", "Measured", "Result"]
    pcm_rows = [pcm_header]
    pcm_data = [
        ["Vth (NMOS)", "0.28", "0.38", "V", "0.33", "PASS"],
        ["Vth (PMOS)", "-0.38", "-0.28", "V", "-0.34", "PASS"],
        ["Idsat (NMOS)", "750", "-", "uA/um", "820", "PASS"],
        ["Idsat (PMOS)", "350", "-", "uA/um", "392", "PASS"],
        ["Ioff (NMOS)", "-", "10", "nA/um", "4.2", "PASS"],
        ["Ioff (PMOS)", "-", "8", "nA/um", "3.1", "PASS"],
        ["Tox", "1.8", "2.2", "nm", "2.01", "PASS"],
        ["Rsh (Poly)", "180", "220", "ohm/sq", "198", "PASS"],
        ["Rsh (Metal1)", "0.04", "0.06", "ohm/sq", "0.049", "PASS"],
    ]
    pcm_rows.extend(pcm_data)
    pcm_table = Table(pcm_rows, colWidths=[28*mm, 22*mm, 22*mm, 18*mm, 22*mm, 16*mm])
    pcm_table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), JP_FONT),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1A237E")),
        ("TEXTCOLOR", (0, 0), (-1, 0), white),
        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    elements.append(pcm_table)
    elements.append(Spacer(1, 3 * mm))
    elements.append(Paragraph(
        "All PCM parameters are within specification limits. No anomalies detected.",
        styles["BodyJP"],
    ))

    elements.append(PageBreak())

    # ─── 4. Appearance Inspection ───
    elements.append(Paragraph("4. Appearance Inspection", styles["HeadingJP"]))
    vis_header = ["Defect Type", "Count", "Spec Limit", "Result"]
    vis_rows = [vis_header,
        ["Scratch", "2", "<=5", "PASS"],
        ["Particle (>0.5um)", "12", "<=30", "PASS"],
        ["Pattern defect", "0", "<=3", "PASS"],
        ["Edge chip", "1", "<=3", "PASS"],
        ["Contamination", "0", "<=2", "PASS"],
    ]
    vis_table = Table(vis_rows, colWidths=[40*mm, 25*mm, 25*mm, 20*mm])
    vis_table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), JP_FONT),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("ALIGN", (1, 0), (-1, -1), "CENTER"),
        ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1A237E")),
        ("TEXTCOLOR", (0, 0), (-1, 0), white),
        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    elements.append(vis_table)

    # ─── 5. Reliability Test ───
    elements.append(Paragraph("5. Reliability Test Results", styles["HeadingJP"]))
    rel_header = ["Test Item", "Conditions", "Duration", "Sample Size", "Failures", "Result"]
    rel_rows = [rel_header,
        ["HTOL", "125C, 1.1xVdd", "1000h", "77 pcs", "0", "PASS"],
        ["TC", "-65C~150C", "1000 cycles", "77 pcs", "0", "PASS"],
        ["HAST", "130C, 85%RH, 1.1xVdd", "96h", "77 pcs", "0", "PASS"],
        ["ESD (HBM)", "Class 2 (>=2kV)", "-", "3 pcs", "0", "PASS"],
        ["ESD (CDM)", "Class C4 (>=500V)", "-", "3 pcs", "0", "PASS"],
        ["Latch-up", "Isup=100mA", "-", "3 pcs", "0", "PASS"],
    ]
    rel_table = Table(rel_rows, colWidths=[22*mm, 35*mm, 22*mm, 20*mm, 18*mm, 16*mm])
    rel_table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), JP_FONT),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1A237E")),
        ("TEXTCOLOR", (0, 0), (-1, 0), white),
        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    elements.append(rel_table)
    elements.append(Spacer(1, 3 * mm))
    elements.append(Paragraph(
        "All reliability tests passed with zero failures. "
        "Product meets JEDEC qualification standards.",
        styles["BodyJP"],
    ))

    # ─── 6. Overall Judgment ───
    elements.append(Paragraph("6. Overall Judgment", styles["HeadingJP"]))
    elements.append(HRFlowable(width="100%", thickness=1, color=HexColor("#4CAF50")))
    elements.append(Spacer(1, 3 * mm))

    judgment_data = [
        ["Overall Result", "PASS"],
        ["Wafer Test", "PASS (Yield: 93.2%)"],
        ["PCM", "PASS (All parameters within spec)"],
        ["Appearance", "PASS (Defects within tolerance)"],
        ["Reliability", "PASS (0 failures across all tests)"],
    ]
    jt = Table(judgment_data, colWidths=[45*mm, 90*mm])
    jt.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), JP_FONT),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("TEXTCOLOR", (0, 0), (0, -1), HexColor("#1A237E")),
        ("TEXTCOLOR", (1, 0), (1, 0), HexColor("#4CAF50")),
        ("BACKGROUND", (0, 0), (-1, 0), HexColor("#E8F5E9")),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LINEBELOW", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
    ]))
    elements.append(jt)

    elements.append(Spacer(1, 8 * mm))
    elements.append(Paragraph(
        "Approved: Quality Assurance Dept. Manager   Signature: _______________   Date: 2026-02-16",
        styles["SmallJP"],
    ))

    doc.build(elements)
    print(f"[OK] PDF Quality Report: {out_path}")


# ═══════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 60)
    print("Semiconductor Sample Input File Generation (SemiDocAI4SPS)")
    print("=" * 60)

    generate_wavedrom_image()
    generate_word_document()
    generate_excel_document()
    generate_pptx_document()
    generate_pdf_document()

    print()
    print("=" * 60)
    print(f"All files generated in {OUTPUT_DIR}")
    print("=" * 60)

    for f in sorted(OUTPUT_DIR.iterdir()):
        size = f.stat().st_size
        print(f"  {f.name:50s} {size:>10,} bytes")
